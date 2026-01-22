import os
import uuid
import sys
import traceback
import logging
from pathlib import Path
from io import BytesIO
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse
from PIL import Image

# Initialize logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Optional libraries with safety checks for Cloud (Vercel/Linux)
try:                                                          
    import fitz  # pymupdf
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False
    logger.warning("pymupdf not found - PDF functionality limited")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx not found - PPTX functionality limited")

# --- CONSOLIDATED CONVERTER LOGIC ---
class FileConverter:
    def __init__(self, upload_dir, download_dir):
        self.upload_dir = Path(upload_dir)
        self.download_dir = Path(download_dir)

    def process_conversion(self, input_filename, target_format):
        input_path = self.upload_dir / input_filename
        input_ext = input_path.suffix.lower()
        output_filename = f"{input_path.stem}.{target_format.lower()}"
        output_path = self.download_dir / output_filename

        try:
            # 1. Image Conversions (JPEG, PNG, etc.)
            if input_ext in [".png", ".jpg", ".jpeg"]:
                img = Image.open(input_path).convert("RGB")
                if target_format.lower() == "pdf":
                    img.save(output_path, "PDF")
                elif target_format.lower() == "pptx" and HAS_PPTX:
                    prs = Presentation()
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide.shapes.add_picture(str(input_path), 0, 0, prs.slide_width, prs.slide_height)
                    prs.save(output_path)
                else:
                    pill_fmt = "JPEG" if target_format.lower() in ["jpg", "jpeg"] else target_format.upper()
                    img.save(output_path, pill_fmt)
                return output_path

            # 2. PDF Conversions (Requires pymupdf)
            elif input_ext == ".pdf" and HAS_FITZ:
                doc = fitz.open(input_path)
                if target_format.lower() in ["png", "jpg", "jpeg"]:
                    page = doc.load_page(0)
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    pix.save(str(output_path))
                    doc.close()
                    return output_path
                elif target_format.lower() == "pptx" and HAS_PPTX:
                    prs = Presentation()
                    for page in doc:
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        img_stream = BytesIO(pix.tobytes("png"))
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        slide.shapes.add_picture(img_stream, 0, 0, prs.slide_width, prs.slide_height)
                    prs.save(output_path)
                    doc.close()
                    return output_path

            raise ValueError(f"Conversion from {input_ext} to {target_format} is not supported on cloud. Try Image or PDF files!")
        except Exception as e:
            logger.error(f"Conversion Error: {e}")
            raise e

# --- VERCEL FASTAPI APP ---
app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Vercel Environment Configuration
IS_CLOUD = os.environ.get('VERCEL') == '1'
TMP_BASE = Path("/tmp") if IS_CLOUD else Path(".")
UPLOAD_DIR = TMP_BASE / "uploads"
DOWNLOAD_DIR = TMP_BASE / "downloads"

UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
DOWNLOAD_DIR.mkdir(parents=True, exist_ok=True)

converter = FileConverter(str(UPLOAD_DIR), str(DOWNLOAD_DIR))

# Simplified task storage (resets on serverless restart)
# In this version, we convert synchronously to avoid polling issues
tasks = {}

@app.exception_handler(Exception)
async def global_exception_handler(request, exc):
    return JSONResponse(
        status_code=500,
        content={"error": "Server Error", "detail": str(exc), "traceback": traceback.format_exc()}
    )

@app.post("/api/upload")
@app.post("/upload")
async def upload_file(file: UploadFile = File(...), target_format: str = "pdf"):
    # Enforce Vercel Hobby 4MB limit
    content = await file.read()
    if len(content) > 4 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File too large (Max 4MB on Vercel Free Plan)")
    
    file_id = str(uuid.uuid4())
    input_filename = f"{file_id}_{file.filename}"
    input_path = UPLOAD_DIR / input_filename
    
    with open(input_path, "wb") as f:
        f.write(content)
    
    # Process conversion immediately for higher reliability on Vercel
    try:
        output_path = converter.process_conversion(input_filename, target_format)
        tasks[file_id] = {
            "status": "completed",
            "output_file": output_path.name,
            "target_format": target_format
        }
        return {"task_id": file_id, "status": "completed"}
    except Exception as e:
        tasks[file_id] = {"status": "failed", "error": str(e)}
        return JSONResponse(status_code=500, content={"error": "Conversion failed", "detail": str(e)})

@app.get("/api/status/{task_id}")
@app.get("/status/{task_id}")
async def get_status(task_id: str):
    if task_id in tasks:
        return tasks[task_id]
    # If the instance reset but the frontend is polling, return completed as a guess 
    # if the file might still be in /tmp (rare but possible) or just 404
    return {"status": "completed"} 

@app.get("/api/download/{task_id}")
@app.get("/download/{task_id}")
async def download_file(task_id: str):
    # Search for files with the task identifier in the download directory
    files = list(DOWNLOAD_DIR.glob(f"*{task_id}*"))
    if not files:
         # Check if the filename is already in the tasks dict
        if task_id in tasks:
            file_path = DOWNLOAD_DIR / tasks[task_id]["output_file"]
            if file_path.exists():
                return FileResponse(file_path, filename=tasks[task_id]["output_file"])
        raise HTTPException(status_code=404, detail="File not found or session expired")
    
    return FileResponse(files[0], filename=files[0].name)

@app.get("/api/health")
@app.get("/health")
async def health():
    return {"status": "ok", "cloud": IS_CLOUD, "fitz": HAS_FITZ, "pptx": HAS_PPTX}

@app.get("/")
@app.get("/api")
async def root():
    return {"message": "EasyConverter Unified API is running!"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
