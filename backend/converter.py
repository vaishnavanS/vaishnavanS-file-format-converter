import os
import shutil
from pathlib import Path
try:
    from pdf2docx import Converter as PDF2DocxConverter
    HAS_PDF2DOCX = True
except ImportError:
    HAS_PDF2DOCX = False

from docx2pdf import convert as docx2pdf_convert
from pptx import Presentation
import fitz # PyMuPDF
from PIL import Image
try:
    import comtypes.client
    HAS_COMTYPES = True
except ImportError:
    HAS_COMTYPES = False

class FileConverter:
    def __init__(self, upload_dir="uploads", download_dir="downloads"):
        self.upload_dir = Path(upload_dir)
        self.download_dir = Path(download_dir)
        self.upload_dir.mkdir(exist_ok=True)
        self.download_dir.mkdir(exist_ok=True)

    def convert_pdf_to_docx(self, input_path, output_path):
        if not HAS_PDF2DOCX:
            raise ImportError("PDF to DOCX conversion is not supported on this platform. Please run locally or use an alternative.")
        cv = PDF2DocxConverter(str(input_path))
        cv.convert(str(output_path))
        cv.close()
        return output_path

    def convert_docx_to_pdf(self, input_path, output_path):
        # On Windows, docx2pdf uses Microsoft Word.
        if os.name != 'nt':
             raise ImportError("DOCX to PDF conversion (High Quality) is not supported on this platform. Please run locally or use an alternative.")
        docx2pdf_convert(str(input_path), str(output_path))
        return output_path

    def convert_pptx_to_pdf(self, input_path, output_path):
        # Uses PowerPoint via COM on Windows
        if not HAS_COMTYPES:
            raise ImportError("PPTX to PDF conversion (High Quality) is not supported on this platform. Please run locally or use an alternative.")
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        # powerpoint.Visible = 1 # Keep it hidden
        try:
            deck = powerpoint.Presentations.Open(str(input_path))
            # 32 is the constant for PDF export
            deck.SaveAs(str(output_path), 32)
            deck.Close()
        finally:
            powerpoint.Quit()
        return output_path

    def convert_pdf_to_pptx(self, input_path, output_path):
        doc = fitz.open(input_path)
        prs = Presentation()
        # Set higher resolution (e.g., 300 DPI)
        zoom = 300 / 72
        matrix = fitz.Matrix(zoom, zoom)
        
        for page in doc:
            pix = page.get_pixmap(matrix=matrix)
            from io import BytesIO
            img_stream = BytesIO(pix.tobytes("png"))
            
            slide_layout = prs.slide_layouts[6] # Blank
            slide = prs.slides.add_slide(slide_layout)
            # Use original dimensions to preserve aspect ratio
            # slide.shapes.add_picture(img_stream, 0, 0, prs.slide_width, prs.slide_height)
            # Better: let it fit the slide effectively
            slide.shapes.add_picture(img_stream, 0, 0, prs.slide_width, prs.slide_height)
            
        prs.save(output_path)
        doc.close()
        return output_path

    def convert_image_to_pptx(self, input_paths, output_path):
        prs = Presentation()
        for path in input_paths:
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(str(path), 0, 0, prs.slide_width, prs.slide_height)
        prs.save(output_path)
        return output_path

    def convert_pdf_to_images(self, input_path, output_dir, format="png"):
        doc = fitz.open(input_path)
        zoom = 300 / 72
        matrix = fitz.Matrix(zoom, zoom)
        image_paths = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=matrix)
            img_path = output_dir / f"page_{i+1}.{format}"
            pix.save(str(img_path))
            image_paths.append(img_path)
        doc.close()
        return image_paths

    def convert_image_to_pdf(self, input_paths, output_path):
        images = []
        for path in input_paths:
            img = Image.open(path).convert("RGB")
            images.append(img)
        
        if images:
            images[0].save(output_path, save_all=True, append_images=images[1:])
        return output_path

    def convert_image_format(self, input_path, output_path, format="PNG"):
        img = Image.open(input_path)
        if format.upper() in ["JPG", "JPEG"]:
            img.convert("RGB").save(output_path, format="JPEG", quality=95, subsampling=0)
        else:
            img.save(output_path, format=format, optimize=True)
        return output_path

    def process_conversion(self, input_filename, target_format):
        input_path = self.upload_dir / input_filename
        input_ext = input_path.suffix.lower()
        output_filename = f"{input_path.stem}.{target_format.lower()}"
        output_path = self.download_dir / output_filename

        try:
            if input_ext == ".pdf":
                if target_format.lower() == "docx":
                    return self.convert_pdf_to_docx(input_path, output_path)
                elif target_format.lower() == "pptx":
                    return self.convert_pdf_to_pptx(input_path, output_path)
                elif target_format.lower() in ["png", "jpg", "jpeg"]:
                    # For images, we might create a zip if there are multiple pages
                    # But for now, let's just do page 1 or return the first path
                    results = self.convert_pdf_to_images(input_path, self.download_dir, target_format.lower())
                    return results[0] # Simplification

            elif input_ext == ".docx":
                if target_format.lower() == "pdf":
                    return self.convert_docx_to_pdf(input_path, output_path)

            elif input_ext in [".pptx", ".ppt"]:
                if target_format.lower() == "pdf":
                    return self.convert_pptx_to_pdf(input_path, output_path)

            elif input_ext in [".png", ".jpg", ".jpeg"]:
                if target_format.lower() == "pdf":
                    return self.convert_image_to_pdf([input_path], output_path)
                elif target_format.lower() == "pptx":
                    return self.convert_image_to_pptx([input_path], output_path)
                else:
                    return self.convert_image_format(input_path, output_path, target_format.upper())

            raise ValueError(f"Unsupported conversion: {input_ext} to {target_format}")
        except Exception as e:
            print(f"Error during conversion: {e}")
            raise e
